"""
文件处理模块

本模块负责处理各种格式的文件上传，提取其中的文本内容。
支持多种常见文档格式，包括文本、PDF、Word、Excel、PPT等。

支持的文件格式：
- 文本文件：.txt, .md
- PDF文档：.pdf
- Word文档：.doc, .docx
- Excel表格：.xls, .xlsx, .csv
- PPT演示：.ppt, .pptx

处理流程：
1. 根据文件扩展名选择对应的处理器
2. 读取文件内容
3. 提取纯文本
4. 返回文本内容供后续处理

依赖库：
- pypdf: PDF处理
- python-docx: Word文档处理
- openpyxl: Excel处理
- python-pptx: PPT处理

注意事项：
- 部分文档格式需要额外安装依赖
- 中文文档需要正确处理编码
- 大文件可能需要分块处理
"""

import os
from typing import Optional
import logging

# 配置日志记录器
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class FileProcessor:
    """
    文件处理器类
    
    提供统一的文件内容提取接口，支持多种文档格式。
    自动根据文件扩展名选择对应的解析器。
    
    Methods:
        process: 主处理方法，根据文件类型自动分发
        _process_text: 处理文本文件
        _process_pdf: 处理PDF文件
        _process_word: 处理Word文档
        _process_excel: 处理Excel文件
        _process_ppt: 处理PPT文件
        _process_csv: 处理CSV文件
    """
    
    async def process(self, file_path: str, file_extension: str) -> Optional[str]:
        """根据文件类型处理文件"""
        ext = file_extension.lower()
        
        if ext == '.txt' or ext == '.md':
            return await self._process_text(file_path)
        elif ext == '.pdf':
            return await self._process_pdf(file_path)
        elif ext in ['.doc', '.docx']:
            return await self._process_word(file_path)
        elif ext in ['.xls', '.xlsx']:
            return await self._process_excel(file_path)
        elif ext in ['.ppt', '.pptx']:
            return await self._process_ppt(file_path)
        elif ext == '.csv':
            return await self._process_csv(file_path)
        else:
            logger.warning(f"不支持的文件类型: {ext}")
            return None
    
    async def _process_text(self, file_path: str) -> str:
        """处理文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()
    
    async def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件"""
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"PDF处理失败: {e}")
            return ""
    
    async def _process_word(self, file_path: str) -> str:
        """处理Word文件"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text.append(" | ".join(row_text))
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Word处理失败: {e}")
            return ""
    
    async def _process_excel(self, file_path: str) -> str:
        """处理Excel文件"""
        try:
            import pandas as pd
            
            # 读取所有sheet
            excel_file = pd.ExcelFile(file_path)
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(df.to_string(index=False))
                text_parts.append("\n")
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Excel处理失败: {e}")
            return ""
    
    async def _process_ppt(self, file_path: str) -> str:
        """处理PPT文件"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                text.append("\n")
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"PPT处理失败: {e}")
            return ""
    
    async def _process_csv(self, file_path: str) -> str:
        """处理CSV文件"""
        try:
            import pandas as pd
            
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"CSV处理失败: {e}")
            return ""